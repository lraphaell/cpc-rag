#!/usr/bin/env python3
"""
Office Document Parser

Extracts text and metadata from Microsoft Office files.
- DOCX: python-docx (paragraphs, headings, tables)
- PPTX: python-pptx (slides with rich text per slide)

For PPTX, generates one rich text chunk per slide by combining:
- Shape text + speaker notes (always)
- Gemini Vision description (for slides with <50 chars of text)
Images are saved to disk for display but NOT embedded as vectors.

Usage:
    from tools.parsers import office_parser
    result = office_parser.parse('/path/to/file.docx')
"""

import io
import os
import subprocess
import tempfile
import time
from pathlib import Path
from typing import Dict, Any, List, Optional


def parse_docx(file_path: str) -> Dict[str, Any]:
    """Parse Word document (.docx)."""
    try:
        from docx import Document

        doc = Document(file_path)

        # Extract paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Extract headings
        headings = []
        for para in doc.paragraphs:
            if para.style.name.startswith('Heading'):
                headings.append({
                    'level': para.style.name,
                    'text': para.text
                })

        # Extract tables
        tables_text = []
        for table_idx, table in enumerate(doc.tables):
            table_content = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    table_content.append(' | '.join(row_text))

            if table_content:
                tables_text.append(f"=== Table {table_idx + 1} ===\n" + '\n'.join(table_content))

        # Combine all text: paragraphs first, then tables
        all_text_parts = []
        if paragraphs:
            all_text_parts.append('\n\n'.join(paragraphs))
        if tables_text:
            all_text_parts.extend(tables_text)

        text = '\n\n'.join(all_text_parts)

        metadata = {
            'parser': 'python-docx',
            'type': 'DOCX',
            'paragraph_count': len(paragraphs),
            'heading_count': len(headings),
            'table_count': len(doc.tables),
            'headings': headings
        }

        return {
            'text': text,
            'metadata': metadata,
            'raw_data': doc
        }

    except ImportError:
        raise ImportError("python-docx not installed. Install with: pip install python-docx")
    except Exception as e:
        raise Exception(f"DOCX parsing failed: {e}")


def _extract_embedded_images(slide, slide_num: int) -> List[Dict]:
    """Extract embedded images (pictures) from a slide."""
    from pptx.shapes.picture import Picture

    images = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            try:
                image = shape.image
                blob = image.blob
                content_type = image.content_type or "image/png"

                # Normalize MIME type
                mime_map = {
                    "image/jpeg": "image/jpeg",
                    "image/jpg": "image/jpeg",
                    "image/png": "image/png",
                    "image/gif": "image/png",  # Gemini prefers png/jpeg
                    "image/bmp": "image/png",
                    "image/tiff": "image/png",
                }
                mime_type = mime_map.get(content_type, "image/png")

                # Convert non-standard formats to PNG via Pillow
                if content_type not in ("image/png", "image/jpeg", "image/jpg"):
                    try:
                        from PIL import Image
                        img = Image.open(io.BytesIO(blob))
                        buf = io.BytesIO()
                        img.save(buf, format="PNG")
                        blob = buf.getvalue()
                        mime_type = "image/png"
                    except Exception:
                        pass  # Use original blob

                images.append({
                    "slide_number": slide_num,
                    "image_bytes": blob,
                    "mime_type": mime_type,
                    "shape_name": shape.name,
                })
            except Exception:
                continue
    return images


def _render_slides_libreoffice(file_path: str) -> Optional[List[Dict]]:
    """
    Render all slides as PNG images using LibreOffice headless.
    Returns None if LibreOffice is not available.
    """
    # Find LibreOffice binary
    lo_paths = [
        "libreoffice",
        "soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/libreoffice",
    ]
    lo_bin = None
    for path in lo_paths:
        try:
            subprocess.run([path, "--version"], capture_output=True, timeout=5)
            lo_bin = path
            break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    if lo_bin is None:
        return None

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            subprocess.run(
                [lo_bin, "--headless", "--convert-to", "png", "--outdir", tmpdir, str(file_path)],
                capture_output=True, timeout=120
            )
            images = []
            for idx, png_path in enumerate(sorted(Path(tmpdir).glob("*.png")), 1):
                images.append({
                    "slide_number": idx,
                    "image_bytes": png_path.read_bytes(),
                    "mime_type": "image/png",
                    "shape_name": f"rendered_slide_{idx}",
                })
            return images if images else None
        except Exception:
            return None


MIN_TEXT_FOR_VISION = 50  # Slides with less text than this get Gemini Vision description


def _describe_slide_with_vision(image_bytes: bytes, slide_text: str, mime_type: str = "image/png") -> str:
    """
    Use Gemini Vision to describe a slide image in text.
    Only called for slides with very little text (<50 chars).
    Returns descriptive text or empty string on failure.
    """
    try:
        from google import genai
        from google.genai import types
        from tools.common.config import GEMINI_API_KEY

        if not GEMINI_API_KEY:
            return ""

        client = genai.Client(api_key=GEMINI_API_KEY)

        prompt = (
            "Describe the content of this presentation slide in detail. "
            "Include all visible text, data from charts/tables, diagram descriptions, "
            "and key information. Write in the same language as the slide content. "
            "Be concise but comprehensive. Focus on factual content, not design."
        )
        if slide_text.strip():
            prompt += f"\n\nText already extracted from this slide: {slide_text.strip()}"

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                types.Part.from_text(text=prompt),
                types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
            ],
        )
        description = response.text.strip() if response.text else ""
        time.sleep(0.5)  # Rate limit protection
        return description

    except Exception as e:
        print(f"      Vision API error: {e}")
        return ""


def parse_pptx(file_path: str) -> Dict[str, Any]:
    """
    Parse PowerPoint presentation (.pptx).

    Strategy: Generate 1 rich text chunk per slide by combining:
    - Shape text + speaker notes (always extracted)
    - Gemini Vision description (for slides with <50 chars of text)

    Images are saved to disk for display but NOT used as vector embeddings.
    This produces much better retrieval than image-only embeddings.
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)

        slide_titles = []
        per_slide_data = []  # One entry per slide with rich text
        vision_count = 0

        # First pass: extract text from all slides
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            title = ""

            # Extract slide title
            if slide.shapes.title:
                title = slide.shapes.title.text
                slide_titles.append(title)

            # Extract all text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

            # Extract speaker notes
            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_content.append(f"[Notas: {notes}]")

            raw_text = '\n'.join(slide_content)

            per_slide_data.append({
                "slide_number": slide_num,
                "title": title,
                "raw_text": raw_text,
                "notes": notes,
                "text_length": len(raw_text.strip()),
            })

        # Extract embedded images for slides that need Vision description
        # and for disk storage (display in frontend)
        slide_images_for_vision = {}
        embedded_images = []
        for slide_num, slide in enumerate(prs.slides, 1):
            images = _extract_embedded_images(slide, slide_num)
            if images:
                embedded_images.extend(images)
                # Keep the largest image per slide for Vision (most likely content, not logo)
                largest = max(images, key=lambda x: len(x.get("image_bytes", b"")))
                slide_images_for_vision[slide_num] = largest

        # Second pass: enrich slides with little text using Gemini Vision
        low_text_slides = [s for s in per_slide_data if s["text_length"] < MIN_TEXT_FOR_VISION]
        if low_text_slides and slide_images_for_vision:
            print(f"    {len(low_text_slides)} slides with <{MIN_TEXT_FOR_VISION} chars text — using Gemini Vision")
            for slide_data in low_text_slides:
                sn = slide_data["slide_number"]
                img_data = slide_images_for_vision.get(sn)
                if img_data:
                    description = _describe_slide_with_vision(
                        img_data["image_bytes"],
                        slide_data["raw_text"],
                        img_data.get("mime_type", "image/png"),
                    )
                    if description:
                        slide_data["raw_text"] += f"\n[Descripción visual: {description}]"
                        slide_data["vision_enriched"] = True
                        vision_count += 1

        # Build final text: one rich chunk per slide
        slides_text = []
        for sd in per_slide_data:
            header = f"=== Slide {sd['slide_number']}"
            if sd["title"]:
                header += f": {sd['title']}"
            header += " ===\n"
            slides_text.append(header + sd["raw_text"])

        text = '\n\n'.join(slides_text)

        if vision_count > 0:
            print(f"    Vision-enriched {vision_count} slides")

        metadata = {
            'parser': 'python-pptx',
            'type': 'PPTX',
            'slide_count': len(prs.slides),
            'slide_titles': slide_titles,
            'image_count': len(embedded_images),
            'vision_enriched_slides': vision_count,
        }

        return {
            'text': text,
            'metadata': metadata,
            'per_slide_data': per_slide_data,  # For per-slide chunking in prepare_chunks
            'slide_images': embedded_images,   # For disk storage (display only, NOT embedding)
            'raw_data': prs,
        }

    except ImportError:
        raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
    except Exception as e:
        raise Exception(f"PPTX parsing failed: {e}")


def parse(file_path: str) -> Dict[str, Any]:
    """
    Parse Office document (auto-detect type).

    Args:
        file_path: Path to DOCX or PPTX file

    Returns:
        dict: {
            'text': str,
            'metadata': dict,
            'raw_data': Any
        }
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    suffix = path.suffix.lower()

    if suffix == '.docx':
        return parse_docx(file_path)
    elif suffix == '.pptx':
        return parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}. Expected .docx or .pptx")


if __name__ == "__main__":
    import argparse

    parser_cli = argparse.ArgumentParser(description="Parse Office document")
    parser_cli.add_argument("file_path", help="Path to DOCX or PPTX file")
    args = parser_cli.parse_args()

    result = parse(args.file_path)
    print(f"Parsed {args.file_path}")
    print(f"Type: {result['metadata']['type']}")
    print(f"Text length: {len(result['text'])} characters")
    if result['metadata']['type'] == 'DOCX':
        print(f"Paragraphs: {result['metadata']['paragraph_count']}")
        print(f"Headings: {result['metadata']['heading_count']}")
    else:
        print(f"Slides: {result['metadata']['slide_count']}")
